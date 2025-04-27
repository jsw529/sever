from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import StreamingResponse, JSONResponse
from google.cloud import texttospeech
from pptx import Presentation
import io

app = FastAPI()

# 여기에 쓰는 거 맞습니다 ✅
client = texttospeech.TextToSpeechClient.from_service_account_file(
    "universal-trail-457914-d8-6d5d654907ed.json"
)

@app.get("/voices")
async def get_voices():
    response = client.list_voices()
    voices_info = [{"name": v.name, "language_codes": v.language_codes} for v in response.voices]
    return voices_info

@app.post("/upload_ppt")
async def upload_ppt(file: UploadFile = File(...)):
    content = await file.read()
    prs = Presentation(io.BytesIO(content))
    all_texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_texts.append(shape.text)
    return {"text": "\\n".join(all_texts)}

@app.post("/synthesize")
async def synthesize(text: str = Form(...), voice_name: str = Form(...)):
    input_text = texttospeech.SynthesisInput(text=text)
    voice = texttospeech.VoiceSelectionParams(
        language_code="ko-KR",
        name=voice_name,
    )
    audio_config = texttospeech.AudioConfig(
        audio_encoding=texttospeech.AudioEncoding.MP3,
        speaking_rate=1.0,
    )

    response = client.synthesize_speech(input=input_text, voice=voice, audio_config=audio_config)
    audio_stream = io.BytesIO(response.audio_content)
    return StreamingResponse(audio_stream, media_type="audio/mpeg")
